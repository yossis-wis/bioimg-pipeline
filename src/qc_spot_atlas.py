"""Spot-atlas QC generation (PowerPoint).

This module reproduces the MATLAB "spot atlas" QC style:
- 15 spots per slide
- multiple zoom levels
- fixed vs adaptive contrast
- tile border colors binned by u0 (spot intensity)
- a yellow outline marking the in-mask used to compute u0

Design
------
This is not part of the strict *contracts* (spots.parquet + manifests).
It's a higher-level QC artifact intended for batch-scale trust building.

Implementation notes
--------------------
We generate each slide as a single PNG (via Matplotlib) and insert it into
PowerPoint (via python-pptx). This matches the old MATLAB approach and keeps
PowerPoint creation robust.
"""

from __future__ import annotations

from dataclasses import dataclass, field
from io import BytesIO
from pathlib import Path
from typing import Any, Dict, Iterable, Optional, Tuple

import numpy as np
import pandas as pd
import yaml

# Matplotlib
# NOTE: We intentionally do NOT force a backend here.
# CLI drivers (e.g. drivers/generate_spot_atlas_pptx.py) set MPLBACKEND=Agg so
# this module works headlessly, while notebooks keep an interactive backend.
import matplotlib.pyplot as plt
from matplotlib.gridspec import GridSpec
from matplotlib.patches import Circle, Rectangle

from image_io import PlaneSelection, read_image_2d

# Keep atlas defaults aligned with Slice0Params (spot detection) when available.
try:
    from slice0_kernel import Slice0Params  # type: ignore

    _SLICE0_DEFAULTS = Slice0Params()
    _DEFAULT_WINDOW_RADIUS_PX = int(_SLICE0_DEFAULTS.window_radius_px)
    _DEFAULT_APERTURE_RADIUS_PX = int(_SLICE0_DEFAULTS.in5_radius_px)
except Exception:  # pragma: no cover
    _DEFAULT_WINDOW_RADIUS_PX = 15
    _DEFAULT_APERTURE_RADIUS_PX = 2


@dataclass(frozen=True)
class SpotAtlasParams:
    """Parameters controlling the spot atlas QC layout."""

    # How crops are extracted (should match Slice0 windowing for interpretability)
    window_radius_px: int = _DEFAULT_WINDOW_RADIUS_PX  # -> (2*window_radius_px+1)x(2*window_radius_px+1)

    # Measurement aperture (yellow outline): should match Slice0Params.in5_radius_px
    aperture_radius_px: int = _DEFAULT_APERTURE_RADIUS_PX

    # Per-slide layout
    spots_per_slide: int = 15

    # Zooms defined in MATLAB as xlim/ylim on a 31x31 crop
    # MATLAB uses 1-based inclusive ranges: [8..25] and [11..21]
    # With 0-based python slicing, those become [7:25] and [10:21]
    mid_slice: slice = field(default_factory=lambda: slice(7, 25))  # 18x18
    tight_slice: slice = field(default_factory=lambda: slice(10, 21))  # 11x11

    # Contrast
    fixed_clim: Optional[Tuple[float, float]] = None  # e.g. (125, 175)
    fixed_percentiles: Tuple[float, float] = (1.0, 99.5)  # used when fixed_clim is None
    nuclei_percentiles: Tuple[float, float] = (1.0, 99.8)

    # Adaptive contrast modeled after the MATLAB deck:
    # clim([bkg, u0 + 1.3*bkg])
    adaptive_bkg_factor: float = 1.3

    # Filtering + ordering
    #
    # u0 (intensity) filtering
    # -----------------------
    #
    # In the integrated workflow, u0 is the Slice0 photometry score:
    #   u0 = mean(in5) - median(out0)
    #
    # This atlas supports *per-channel* filtering because DNA-locus spots and
    # protein spots can have very different intensity distributions.
    #
    # - u0_min is the default threshold (applies when a channel has no override).
    # - u0_min_by_channel overrides the default for specific spot_channel values.
    #
    u0_min: float = 0.0
    u0_min_by_channel: Dict[int, float] = field(default_factory=dict)

    # If set, restrict the atlas to nuclei that contain at least one spot from
    # *each* of these spot_channel values. This is useful for multi-channel
    # co-localization QC (e.g. DNA + protein spots).
    require_spot_channels: Optional[Tuple[int, ...]] = None

    # Optional: restrict which spot channels are *rendered* as columns in the deck.
    #
    # This is orthogonal to ``require_spot_channels``:
    # - ``require_spot_channels`` filters *nuclei* (co-occurrence requirement).
    # - ``include_spot_channels`` filters which spots become atlas columns.
    #
    # Example (DNA+protein co-occurrence QC):
    #   - require_spot_channels = (1, 2)   (nuclei must contain both)
    #   - include_spot_channels = (2,)     (show only protein spots as columns)
    include_spot_channels: Optional[Tuple[int, ...]] = None

    # Optional top context row: show a zoomed-out crop from a reference spot channel
    # (typically the DNA-locus channel) and mark the reference spot location.
    #
    # If ``context_spot_channel`` is None, the context row is omitted.
    context_spot_channel: Optional[int] = None
    context_window_radius_px: int = 30
    context_marker_radius_px: int = 4

    sort_by: str = "intensity"  # intensity | snr | random | input_path
    random_seed: int = 0

    # Rendering
    figure_size_in: Tuple[float, float] = (13.333, 7.5)  # 16:9
    dpi: int = 150
    title_fontsize: int = 18
    legend_fontsize: int = 14

    # Performance
    image_cache_items: int = 8


# --- u0 bin colors (match MATLAB legend) ---
_U0_BINS = [15, 20, 25, 30, 35, 40, 45, 50, 55]
_U0_COLORS = [
    (1.0, 0.0, 0.0),  # 15-20 red
    (0.0, 1.0, 0.0),  # 20-25 green
    (0.0, 0.0, 1.0),  # 25-30 blue
    (0.0, 1.0, 1.0),  # 30-35 cyan
    (1.0, 0.0, 1.0),  # 35-40 magenta
    (1.0, 1.0, 0.0),  # 40-45 yellow
    (0.0, 0.0, 0.0),  # 45-50 black
    (1.0, 1.0, 1.0),  # 50-55 white
    (0.7, 0.7, 0.7),  # >55 gray
]
_U0_LABELS = [
    "15-20",
    "20-25",
    "25-30",
    "30-35",
    "35-40",
    "40-45",
    "45-50",
    "50-55",
    ">55",
]


def _u0_to_color(u0: float) -> Tuple[float, float, float]:
    u0 = float(u0)
    if u0 < 20:
        return _U0_COLORS[0]
    if u0 < 25:
        return _U0_COLORS[1]
    if u0 < 30:
        return _U0_COLORS[2]
    if u0 < 35:
        return _U0_COLORS[3]
    if u0 < 40:
        return _U0_COLORS[4]
    if u0 < 45:
        return _U0_COLORS[5]
    if u0 < 50:
        return _U0_COLORS[6]
    if u0 < 55:
        return _U0_COLORS[7]
    return _U0_COLORS[8]


def _robust_clim(img: np.ndarray, pmin: float, pmax: float) -> Tuple[float, float]:
    x = np.asarray(img, dtype=float)
    if x.size == 0:
        return 0.0, 1.0
    lo, hi = np.percentile(x, [float(pmin), float(pmax)])
    if not np.isfinite(lo) or not np.isfinite(hi) or hi <= lo:
        lo = float(np.min(x))
        hi = float(np.max(x))
        if not np.isfinite(lo) or not np.isfinite(hi) or hi <= lo:
            return 0.0, 1.0
    return float(lo), float(hi)


def _extract_centered_window(img: np.ndarray, y: int, x: int, radius: int) -> np.ndarray:
    """Extract (2r+1)x(2r+1) centered window.

    This function assumes the window is fully inside the image.
    Slice0 detection already enforces that; we keep this strict so QC matches detection.
    """

    r = int(radius)
    y0 = int(y) - r
    y1 = int(y) + r + 1
    x0 = int(x) - r
    x1 = int(x) + r + 1
    return img[y0:y1, x0:x1]


def _extract_centered_window_with_padding(
    img: np.ndarray,
    y: int,
    x: int,
    radius: int,
    *,
    pad_value: float = 0.0,
) -> np.ndarray:
    """Extract a centered window with zero-padding at the borders.

    Unlike :func:`_extract_centered_window`, this helper is *not* used for
    photometry QC. It's only for the zoomed-out context row, where we prefer to
    show as much neighborhood context as possible even if a spot is near the
    image edge.
    """

    r = int(radius)
    h, w = img.shape[:2]
    out = np.full((2 * r + 1, 2 * r + 1), pad_value, dtype=img.dtype)

    y0 = int(y) - r
    x0 = int(x) - r
    y1 = int(y) + r + 1
    x1 = int(x) + r + 1

    src_y0 = max(y0, 0)
    src_x0 = max(x0, 0)
    src_y1 = min(y1, h)
    src_x1 = min(x1, w)

    dst_y0 = src_y0 - y0
    dst_x0 = src_x0 - x0
    dst_y1 = dst_y0 + (src_y1 - src_y0)
    dst_x1 = dst_x0 + (src_x1 - src_x0)

    if src_y1 > src_y0 and src_x1 > src_x0:
        out[dst_y0:dst_y1, dst_x0:dst_x1] = img[src_y0:src_y1, src_x0:src_x1]

    return out



def _disk_mask(radius_px: int, window_size: int) -> np.ndarray:
    """Binary disk mask centered in a window."""
    from skimage.morphology import disk

    base = disk(int(radius_px)).astype(bool)
    pad = (int(window_size) - int(base.shape[0])) // 2
    if pad < 0:
        raise ValueError("disk mask larger than window")
    return np.pad(base, ((pad, pad), (pad, pad)), mode="constant", constant_values=0).astype(bool)


def _mask_outline_xy(mask: np.ndarray) -> Tuple[np.ndarray, np.ndarray]:
    """Return a MATLAB-style pixel-edge outline (x,y) for a binary mask.

    MATLAB's deck drew the measurement aperture as a polyline that traces the
    *outside edges* of the included pixels (vertices land on pixel corners at
    half-integers, e.g. 13.5, 14.5, ...). This is important for QC because it
    makes it visually unambiguous which pixels are included in the mean.

    Coordinate convention
    ---------------------
    We render crops using:

        extent = (0.5, W+0.5, H+0.5, 0.5)

    so pixel centers are at integer coordinates 1..W (MATLAB-like), and pixel
    corners lie on half-integers. This function returns coordinates in that
    same space (already shifted by +1 relative to 0-based array indices).

    Implementation
    --------------
    We treat each True pixel as a unit square and extract the *external* set of
    boundary edges (edges not shared with another True pixel). Those edges form
    one or more closed loops on the pixel-corner lattice; we traverse the loop
    with the lexicographically-smallest start vertex and return it as a closed
    polyline.
    """
    m = np.asarray(mask, dtype=bool)
    if m.ndim != 2:
        raise ValueError(f"mask must be 2D, got shape={m.shape}")

    h, w = m.shape
    if h == 0 or w == 0 or not np.any(m):
        return np.asarray([]), np.asarray([])

    # Collect external edges as undirected segments between pixel-corner vertices.
    # Vertices are in the MATLAB-like coordinate system:
    #   pixel (i,j) covers x in [j+0.5, j+1.5], y in [i+0.5, i+1.5]
    edges: set[tuple[tuple[float, float], tuple[float, float]]] = set()

    def _add_edge(v1: tuple[float, float], v2: tuple[float, float]) -> None:
        # Store undirected edge with deterministic ordering
        edges.add((v1, v2) if v1 <= v2 else (v2, v1))

    for i in range(h):
        for j in range(w):
            if not m[i, j]:
                continue

            x0 = j + 0.5
            x1 = j + 1.5
            y0 = i + 0.5
            y1 = i + 1.5

            # Top edge (no True neighbor above)
            if i == 0 or not m[i - 1, j]:
                _add_edge((x0, y0), (x1, y0))
            # Bottom edge
            if i == h - 1 or not m[i + 1, j]:
                _add_edge((x0, y1), (x1, y1))
            # Left edge
            if j == 0 or not m[i, j - 1]:
                _add_edge((x0, y0), (x0, y1))
            # Right edge
            if j == w - 1 or not m[i, j + 1]:
                _add_edge((x1, y0), (x1, y1))

    if not edges:
        return np.asarray([]), np.asarray([])

    # Build vertex adjacency (each boundary vertex should have degree 2 for a simple loop).
    adj: dict[tuple[float, float], list[tuple[float, float]]] = {}
    for v1, v2 in edges:
        adj.setdefault(v1, []).append(v2)
        adj.setdefault(v2, []).append(v1)

    # Pick a deterministic start: smallest y, then smallest x
    start = min(adj.keys(), key=lambda v: (v[1], v[0]))

    # Traverse the loop
    path: list[tuple[float, float]] = [start]
    prev: Optional[tuple[float, float]] = None
    curr: tuple[float, float] = start

    for _ in range(len(edges) + 10):
        nbrs = adj.get(curr, [])
        if not nbrs:
            break

        if prev is None:
            # choose a deterministic first direction
            nxt = sorted(nbrs, key=lambda v: (v[1], v[0]))[0]
        else:
            if len(nbrs) == 1:
                nxt = nbrs[0]
            else:
                nxt = nbrs[0] if nbrs[0] != prev else nbrs[1]

        if nxt == start:
            path.append(start)
            break

        path.append(nxt)
        prev, curr = curr, nxt

    xs = np.asarray([p[0] for p in path], dtype=float)
    ys = np.asarray([p[1] for p in path], dtype=float)
    return xs, ys


class _LRUImageCache:
    def __init__(self, max_items: int):
        self.max_items = int(max_items)
        self._store: Dict[Tuple[Any, ...], Tuple[np.ndarray, ...]] = {}
        self._order: list[Tuple[Any, ...]] = []

    def get(self, key: Tuple[Any, ...]) -> Optional[Tuple[np.ndarray, ...]]:
        if key not in self._store:
            return None
        # bump
        try:
            self._order.remove(key)
        except ValueError:
            pass
        self._order.append(key)
        return self._store[key]

    def put(self, key: Tuple[Any, ...], value: Tuple[np.ndarray, ...]) -> None:
        if key in self._store:
            try:
                self._order.remove(key)
            except ValueError:
                pass
        self._store[key] = value
        self._order.append(key)
        while len(self._order) > self.max_items:
            old = self._order.pop(0)
            self._store.pop(old, None)


def load_run_manifest(run_dir: Path) -> Dict[str, Any]:
    path = Path(run_dir) / "run_manifest.yaml"
    with path.open("r", encoding="utf-8") as f:
        data = yaml.safe_load(f) or {}
    if not isinstance(data, dict):
        raise ValueError(f"Invalid run_manifest.yaml: {path}")
    return data


def _plane_selection_from_manifest(manifest: Dict[str, Any], channel: int) -> PlaneSelection:
    cfg = manifest.get("config_snapshot", {})
    if not isinstance(cfg, dict):
        cfg = {}
    return PlaneSelection(
        channel=int(channel),
        ims_resolution_level=int(cfg.get("ims_resolution_level", 0)),
        ims_time_index=int(cfg.get("ims_time_index", 0)),
        ims_z_index=int(cfg.get("ims_z_index", 0)),
    )


def _channel_nuclei_from_manifest(manifest: Dict[str, Any]) -> Optional[int]:
    cfg = manifest.get("config_snapshot", {})
    if not isinstance(cfg, dict):
        return None
    raw = cfg.get("channel_nuclei", None)
    if raw in (None, 0, "none", "None"):
        return None
    return int(raw)


def _resolve_input_path(manifest: Dict[str, Any]) -> Path:
    p = manifest.get("input_path")
    if not p:
        raise ValueError("run_manifest missing input_path")
    return Path(str(p)).expanduser().resolve()


def _infer_title_prefix(spots_df: pd.DataFrame) -> str:
    parts = []
    if "condition" in spots_df.columns:
        uniq = spots_df["condition"].dropna().astype(str).unique().tolist()
        if len(uniq) == 1 and uniq[0]:
            parts.append(uniq[0])
    if "spot_channel" in spots_df.columns:
        chs = sorted({int(c) for c in spots_df["spot_channel"].dropna().tolist()})
        if len(chs) == 1:
            parts.append(f"ch{chs[0]}")
    return " / ".join(parts)


def _sort_spots(df: pd.DataFrame, params: SpotAtlasParams) -> pd.DataFrame:
    mode = str(params.sort_by).lower()
    if mode == "snr" and "snr" in df.columns:
        return df.sort_values("snr", ascending=False, kind="mergesort")
    if mode == "random":
        rng = np.random.default_rng(int(params.random_seed))
        order = rng.permutation(len(df))
        return df.iloc[order]
    if mode == "input_path" and "input_path" in df.columns:
        return df.sort_values(["input_path", "intensity"], ascending=[True, False], kind="mergesort")
    # default: intensity
    return df.sort_values("intensity", ascending=False, kind="mergesort")


def _ensure_required_spots_columns(df: pd.DataFrame) -> None:
    required = {"y_px", "x_px", "intensity", "background", "spot_channel", "output_dir"}
    missing = [c for c in required if c not in df.columns]
    if missing:
        raise ValueError(
            "spots table is missing required columns for atlas generation: " + ", ".join(missing)
        )


def _u0_threshold_for_channel(params: SpotAtlasParams, spot_channel: int) -> float:
    """Return the u0 (intensity) threshold for a given spot channel."""

    ch = int(spot_channel)
    if params.u0_min_by_channel and ch in params.u0_min_by_channel:
        return float(params.u0_min_by_channel[ch])
    return float(params.u0_min)


def _format_u0_filter(params: SpotAtlasParams, channels: Iterable[int]) -> str:
    """Compact string describing per-channel u0 thresholds (for slide titles)."""

    chs = sorted({int(c) for c in channels})
    if not chs:
        return f"u0 >= {float(params.u0_min):g}"

    # If no per-channel overrides are set, keep the old title format.
    if not params.u0_min_by_channel:
        return f"u0 >= {float(params.u0_min):g}"

    parts = []
    for ch in chs:
        thr = _u0_threshold_for_channel(params, ch)
        parts.append(f"ch{ch}>={thr:g}")
    return "u0 " + ",".join(parts)


def _filter_spots_by_u0(df: pd.DataFrame, params: SpotAtlasParams) -> pd.DataFrame:
    """Filter a spot table by per-channel u0 thresholds."""

    if df.empty:
        return df

    out = df.copy()
    out = out[np.isfinite(out["y_px"]) & np.isfinite(out["x_px"])].copy()

    # Compute a per-row threshold from spot_channel.
    try:
        ch_series = out["spot_channel"].astype(int)
    except Exception:
        ch_series = out["spot_channel"].apply(lambda v: int(v))

    thr_map = {int(k): float(v) for k, v in (params.u0_min_by_channel or {}).items()}
    default_thr = float(params.u0_min)
    thr = ch_series.map(thr_map).fillna(default_thr).astype(float)

    u0 = out["intensity"].astype(float)
    out = out[u0 >= thr].copy()
    return out


def _filter_spots_to_required_nuclei(
    df: pd.DataFrame,
    *,
    required_spot_channels: Iterable[int],
) -> pd.DataFrame:
    """Restrict spots to nuclei containing all required spot channels.

    A nucleus identity is defined as (output_dir, nucleus_label) to avoid label
    collisions across different input files.
    """

    required = {int(c) for c in required_spot_channels}
    if not required:
        return df

    needed_cols = {"output_dir", "nucleus_label", "spot_channel"}
    missing = [c for c in needed_cols if c not in df.columns]
    if missing:
        raise ValueError(
            "require_spot_channels needs columns: " + ", ".join(missing)
        )

    if df.empty:
        return df

    out = df.copy()
    out["spot_channel"] = out["spot_channel"].astype(int)
    out["nucleus_label"] = out["nucleus_label"].astype(int)

    # Only consider spots assigned to a nucleus label > 0.
    out = out[out["nucleus_label"] > 0].copy()
    if out.empty:
        return out

    # Compute, for each nucleus identity, which spot channels appear.
    ch_sets = (
        out.groupby(["output_dir", "nucleus_label"], sort=False)["spot_channel"]
        .apply(lambda s: set(s.dropna().astype(int)))
        .rename("spot_channels_present")
    )

    ok = ch_sets.apply(lambda present: required.issubset(present))
    nuclei = ch_sets[ok].reset_index().loc[:, ["output_dir", "nucleus_label"]]
    if nuclei.empty:
        return out.iloc[0:0].copy()

    # Filter spots down to those nuclei.
    return out.merge(nuclei, on=["output_dir", "nucleus_label"], how="inner")


def _best_spot_by_nucleus(
    df: pd.DataFrame,
    *,
    spot_channel: int,
) -> Dict[Tuple[str, int], Tuple[float, float]]:
    """Return (output_dir, nucleus_label) -> (y_px, x_px) for the brightest spot.

    We take the *maximum u0* (stored in the contract column ``intensity``) within
    each nucleus for the given ``spot_channel``. This is useful for atlas QC
    where we want to mark (e.g.) the most likely DNA-locus spot while browsing
    protein spots.

    Notes
    -----
    - Nucleus identity uses ``output_dir`` to avoid label collisions across runs.
    - Nuclei with label <= 0 are ignored.
    """

    if df.empty:
        return {}

    if "nucleus_label" not in df.columns:
        return {}

    out = df.copy()
    out["spot_channel"] = out["spot_channel"].astype(int)
    out["nucleus_label"] = out["nucleus_label"].astype(int)

    out = out[(out["spot_channel"] == int(spot_channel)) & (out["nucleus_label"] > 0)].copy()
    if out.empty:
        return {}

    # For each nucleus, keep the brightest (max intensity) spot.
    idx = out.groupby(["output_dir", "nucleus_label"], sort=False)["intensity"].idxmax()
    best = out.loc[idx, ["output_dir", "nucleus_label", "y_px", "x_px"]].copy()

    mapping: Dict[Tuple[str, int], Tuple[float, float]] = {}
    for _, r in best.iterrows():
        key = (str(Path(str(r["output_dir"]))), int(r["nucleus_label"]))
        mapping[key] = (float(r["y_px"]), float(r["x_px"]))
    return mapping


def render_atlas_page_png(
    spots_batch: pd.DataFrame,
    *,
    page_title: str,
    params: SpotAtlasParams,
    image_cache: _LRUImageCache,
    manifest_cache: Dict[str, Dict[str, Any]],
    fixed_clim_state: Dict[str, Tuple[float, float]],
    context_spots_by_nucleus: Optional[Dict[Tuple[str, int], Tuple[float, float]]] = None,
) -> bytes:
    """Render a single atlas page (one slide image) and return PNG bytes."""

    spots_batch = spots_batch.reset_index(drop=True)

    win_r = int(params.window_radius_px)
    win = 2 * win_r + 1

    # Measurement aperture outline: matches Slice0's disk radius=2 by default.
    in5 = _disk_mask(int(params.aperture_radius_px), win)
    outline_x, outline_y = _mask_outline_xy(in5)

    extent_full = (0.5, win + 0.5, win + 0.5, 0.5)

    # Optional extra context row (top): zoomed-out crop from a reference channel.
    use_context = params.context_spot_channel not in (None, 0, "none", "None")
    ctx_channel = int(params.context_spot_channel) if use_context else 0
    ctx_r = int(params.context_window_radius_px) if use_context else 0
    ctx_win = 2 * ctx_r + 1 if use_context else 0
    extent_ctx = (0.5, ctx_win + 0.5, ctx_win + 0.5, 0.5) if use_context else extent_full

    n_tile_rows = 8 if use_context else 7
    base_row = 1 if use_context else 0
    legend_row = n_tile_rows

    fig = plt.figure(figsize=params.figure_size_in, dpi=int(params.dpi), facecolor="white")
    gs = GridSpec(
        nrows=n_tile_rows + 1,
        ncols=int(params.spots_per_slide),
        figure=fig,
        wspace=0.0,
        hspace=0.0,
        height_ratios=[1] * n_tile_rows + [0.7],
    )

    cmap = "gray_r"  # match MATLAB's flipud(gray)

    for col in range(int(params.spots_per_slide)):
        if col >= len(spots_batch):
            # still create empty axes to keep layout stable
            for rr in range(n_tile_rows):
                ax = fig.add_subplot(gs[rr, col])
                ax.axis("off")
            continue

        row = spots_batch.iloc[col]
        u0 = float(row["intensity"])
        bkg = float(row.get("background", 0.0))
        y_px = int(round(float(row["y_px"])))
        x_px = int(round(float(row["x_px"])))
        nucleus_label = int(row.get("nucleus_label", 0))

        out_dir = Path(str(row["output_dir"]))
        out_key = str(out_dir)
        if out_key not in manifest_cache:
            manifest_cache[out_key] = load_run_manifest(out_dir)
        manifest = manifest_cache[out_key]

        input_path = _resolve_input_path(manifest)
        spot_channel = int(row.get("spot_channel", 1))
        nuc_channel = _channel_nuclei_from_manifest(manifest)

        sel_spot = _plane_selection_from_manifest(manifest, spot_channel)
        sel_nuc = _plane_selection_from_manifest(manifest, nuc_channel) if nuc_channel else None
        sel_ctx = _plane_selection_from_manifest(manifest, ctx_channel) if use_context else None

        cache_key = (
            str(input_path),
            int(spot_channel),
            int(ctx_channel),
            int(sel_spot.ims_resolution_level),
            int(sel_spot.ims_time_index),
            int(sel_spot.ims_z_index),
            int(nuc_channel or 0),
        )
        cached = image_cache.get(cache_key)
        if cached is None:
            img_spot = read_image_2d(input_path, sel_spot)
            img_nuc = read_image_2d(input_path, sel_nuc) if sel_nuc else np.zeros_like(img_spot)

            if use_context:
                if ctx_channel == spot_channel:
                    img_ctx = img_spot
                elif nuc_channel is not None and ctx_channel == int(nuc_channel):
                    img_ctx = img_nuc
                else:
                    if sel_ctx is None:
                        raise ValueError("context_spot_channel set but context plane selection is None")
                    img_ctx = read_image_2d(input_path, sel_ctx)
            else:
                img_ctx = img_spot

            image_cache.put(cache_key, (img_nuc, img_spot, img_ctx))
        else:
            img_nuc, img_spot, img_ctx = cached

        nuc_crop = _extract_centered_window(img_nuc, y_px, x_px, win_r)
        spot_crop = _extract_centered_window(img_spot, y_px, x_px, win_r)

        ctx_crop: Optional[np.ndarray]
        ctx_vmin: float
        ctx_vmax: float
        if use_context:
            ctx_crop = _extract_centered_window_with_padding(img_ctx, y_px, x_px, ctx_r)
            ctx_vmin, ctx_vmax = _robust_clim(ctx_crop, params.fixed_percentiles[0], params.fixed_percentiles[1])
        else:
            ctx_crop = None
            ctx_vmin, ctx_vmax = 0.0, 1.0

        # Fixed clim (global across the deck unless user specified explicit fixed_clim)
        if params.fixed_clim is not None:
            fixed_vmin, fixed_vmax = float(params.fixed_clim[0]), float(params.fixed_clim[1])
        else:
            if "fixed" not in fixed_clim_state:
                fixed_clim_state["fixed"] = _robust_clim(
                    img_spot, params.fixed_percentiles[0], params.fixed_percentiles[1]
                )
            fixed_vmin, fixed_vmax = fixed_clim_state["fixed"]

        # Nuclei scaling per-crop (robust)
        nuc_vmin, nuc_vmax = _robust_clim(nuc_crop, params.nuclei_percentiles[0], params.nuclei_percentiles[1])

        # Adaptive clim (match MATLAB heuristic)
        adapt_vmin = float(bkg)
        adapt_vmax = float(u0) + float(params.adaptive_bkg_factor) * float(bkg)
        if adapt_vmax <= adapt_vmin:
            adapt_vmax = adapt_vmin + 1.0

        # Prepare zooms
        mid = params.mid_slice
        tight = params.tight_slice

        spot_mid = spot_crop[mid, mid]
        spot_tight = spot_crop[tight, tight]

        extent_mid = (0.5, spot_mid.shape[1] + 0.5, spot_mid.shape[0] + 0.5, 0.5)
        extent_tight = (
            0.5,
            spot_tight.shape[1] + 0.5,
            spot_tight.shape[0] + 0.5,
            0.5,
        )

        border_color = _u0_to_color(u0)

        def _style_ax(ax: plt.Axes) -> None:
            ax.set_xticks([])
            ax.set_yticks([])
            for spine in ax.spines.values():
                spine.set_linewidth(2.0)
                spine.set_edgecolor(border_color)

        # Optional Row 0: zoomed-out context crop (reference channel) with DNA spot marker.
        if use_context and ctx_crop is not None:
            ax_ctx = fig.add_subplot(gs[0, col])
            ax_ctx.imshow(
                ctx_crop,
                cmap=cmap,
                vmin=ctx_vmin,
                vmax=ctx_vmax,
                interpolation="nearest",
                extent=extent_ctx,
            )

            if context_spots_by_nucleus is not None and nucleus_label > 0:
                key = (out_key, nucleus_label)
                if key in context_spots_by_nucleus:
                    ref_y, ref_x = context_spots_by_nucleus[key]
                    x_ref = (ctx_r + 1) + (float(ref_x) - float(x_px))
                    y_ref = (ctx_r + 1) + (float(ref_y) - float(y_px))
                    ax_ctx.add_patch(
                        Circle(
                            (x_ref, y_ref),
                            radius=float(params.context_marker_radius_px),
                            fill=False,
                            edgecolor="yellow",
                            linewidth=1.0,
                        )
                    )

            _style_ax(ax_ctx)

        # Row: nuclei (context) + aperture outline
        ax0 = fig.add_subplot(gs[base_row + 0, col])
        ax0.imshow(
            nuc_crop,
            cmap=cmap,
            vmin=nuc_vmin,
            vmax=nuc_vmax,
            interpolation="nearest",
            extent=extent_full,
        )
        if outline_x.size:
            ax0.plot(outline_x, outline_y, color="yellow", linewidth=1.0)
        _style_ax(ax0)

        # Row: full spot (fixed)
        ax1 = fig.add_subplot(gs[base_row + 1, col])
        ax1.imshow(
            spot_crop,
            cmap=cmap,
            vmin=fixed_vmin,
            vmax=fixed_vmax,
            interpolation="nearest",
            extent=extent_full,
        )
        _style_ax(ax1)

        # Row: mid spot (fixed)
        ax2 = fig.add_subplot(gs[base_row + 2, col])
        ax2.imshow(
            spot_mid,
            cmap=cmap,
            vmin=fixed_vmin,
            vmax=fixed_vmax,
            interpolation="nearest",
            extent=extent_mid,
        )
        _style_ax(ax2)

        # Row: tight spot (fixed) + outline
        ax3 = fig.add_subplot(gs[base_row + 3, col])
        ax3.imshow(
            spot_tight,
            cmap=cmap,
            vmin=fixed_vmin,
            vmax=fixed_vmax,
            interpolation="nearest",
            extent=extent_tight,
        )
        if outline_x.size:
            ax3.plot(outline_x - tight.start, outline_y - tight.start, color="yellow", linewidth=1.0)
        _style_ax(ax3)

        # Row: full spot (adaptive)
        ax4 = fig.add_subplot(gs[base_row + 4, col])
        ax4.imshow(
            spot_crop,
            cmap=cmap,
            vmin=adapt_vmin,
            vmax=adapt_vmax,
            interpolation="nearest",
            extent=extent_full,
        )
        _style_ax(ax4)

        # Row: mid spot (adaptive)
        ax5 = fig.add_subplot(gs[base_row + 5, col])
        ax5.imshow(
            spot_mid,
            cmap=cmap,
            vmin=adapt_vmin,
            vmax=adapt_vmax,
            interpolation="nearest",
            extent=extent_mid,
        )
        _style_ax(ax5)

        # Row: tight spot (adaptive) + outline
        ax6 = fig.add_subplot(gs[base_row + 6, col])
        ax6.imshow(
            spot_tight,
            cmap=cmap,
            vmin=adapt_vmin,
            vmax=adapt_vmax,
            interpolation="nearest",
            extent=extent_tight,
        )
        if outline_x.size:
            ax6.plot(outline_x - tight.start, outline_y - tight.start, color="yellow", linewidth=1.0)
        _style_ax(ax6)

    # Legend row spans all columns
    ax_leg = fig.add_subplot(gs[legend_row, :])
    ax_leg.set_xlim(0, 9)
    ax_leg.set_ylim(0, 1)
    ax_leg.axis("off")
    for i, (color, label) in enumerate(zip(_U0_COLORS, _U0_LABELS)):
        ax_leg.add_patch(Rectangle((i, 0), 1, 1, facecolor=color, edgecolor="black", linewidth=0.5))
        ax_leg.text(
            i + 0.5,
            0.5,
            label,
            ha="center",
            va="center",
            fontsize=int(params.legend_fontsize),
            color="black" if sum(color) > 1.5 else "white" if color == (0.0, 0.0, 0.0) else "black",
        )

    fig.suptitle(page_title, fontsize=int(params.title_fontsize), y=0.99)

    bio = BytesIO()
    fig.savefig(bio, format="png", bbox_inches="tight", pad_inches=0.0)
    plt.close(fig)
    return bio.getvalue()


def build_spot_atlas_pptx(
    spots_df: pd.DataFrame,
    *,
    out_pptx: Path,
    params: SpotAtlasParams,
    group_by: Optional[str] = None,
    deck_title: Optional[str] = None,
) -> Path:
    """Build an atlas PPTX from an aggregate spots table.

    Parameters
    ----------
    spots_df:
        Aggregate spots table. Must include: y_px, x_px, intensity, background,
        spot_channel, output_dir. If input_path exists it is used in titles.
    out_pptx:
        Destination PPTX path.
    params:
        SpotAtlasParams.
    group_by:
        Optional grouping column name (e.g. "condition" or "spot_channel").
        Each group becomes a separate section in the deck.
    deck_title:
        Optional prefix for slide titles.

    Returns
    -------
    Path
        out_pptx
    """

    _ensure_required_spots_columns(spots_df)

    try:
        from pptx import Presentation
        from pptx.util import Inches
    except Exception as e:
        raise ImportError(
            "python-pptx is required to generate PowerPoint QC decks. "
            "Install it via environment.yml (conda-forge: python-pptx)."
        ) from e

    out_pptx = Path(out_pptx).expanduser().resolve()
    out_pptx.parent.mkdir(parents=True, exist_ok=True)

    # 1) Photometry (u0) filter: supports global + per-channel thresholds.
    df_all = _filter_spots_by_u0(spots_df, params)

    # 2) Optional nucleus-level co-occurrence filter.
    if params.require_spot_channels:
        df_all = _filter_spots_to_required_nuclei(df_all, required_spot_channels=params.require_spot_channels)

    if df_all.empty:
        chs: list[int] = []
        try:
            chs = sorted({int(c) for c in spots_df["spot_channel"].dropna().unique()})
        except Exception:
            chs = []
        filt_desc = _format_u0_filter(params, chs)
        extra = ""
        if params.require_spot_channels:
            extra = f" and requiring nuclei with channels={tuple(int(c) for c in params.require_spot_channels)}"
        raise ValueError(f"No spots remain after filtering ({filt_desc}){extra}")

    # 3) Context marker mapping (e.g., DNA channel) computed on the *full* filtered set.
    context_spots_by_nucleus: Optional[Dict[Tuple[str, int], Tuple[float, float]]] = None
    if params.context_spot_channel not in (None, 0, "none", "None"):
        context_spots_by_nucleus = _best_spot_by_nucleus(df_all, spot_channel=int(params.context_spot_channel))

    # 4) Optional: choose which spot channels are rendered as atlas columns.
    df = df_all
    include_desc = ""
    if params.include_spot_channels:
        include = sorted({int(c) for c in params.include_spot_channels})
        include_desc = "show " + ",".join(f"ch{c}" for c in include)
        df = df.copy()
        df["spot_channel"] = df["spot_channel"].astype(int)
        df = df[df["spot_channel"].isin(include)].copy()
        if df.empty:
            raise ValueError(f"No spots remain after include_spot_channels filter ({include_desc}).")

    df = _sort_spots(df, params)

    pres = Presentation()
    pres.slide_width = Inches(float(params.figure_size_in[0]))
    pres.slide_height = Inches(float(params.figure_size_in[1]))
    blank = pres.slide_layouts[6]

    image_cache = _LRUImageCache(params.image_cache_items)
    manifest_cache: Dict[str, Dict[str, Any]] = {}
    fixed_clim_state: Dict[str, Tuple[float, float]] = {}

    if group_by and group_by in df.columns:
        grouped = df.groupby(group_by, sort=False)
    else:
        grouped = [(None, df)]

    title_prefix = deck_title or _infer_title_prefix(df)

    for g_name, g_df in grouped:
        g_df = g_df.reset_index(drop=True)

        # MATLAB-like: sort by intensity within group
        g_df = _sort_spots(g_df, params)

        for start in range(0, len(g_df), int(params.spots_per_slide)):
            batch = g_df.iloc[start : start + int(params.spots_per_slide)]

            parts: list[str] = []
            if title_prefix:
                parts.append(title_prefix)
            if g_name is not None:
                parts.append(f"{group_by}={g_name}")
            parts.append(
                f"spots {start + 1}-{min(start + int(params.spots_per_slide), len(g_df))} / {len(g_df)}"
            )

            # For the u0 filter text, prefer the co-occurrence channels (if set) so
            # the title reflects *both* spot types even if we only render one channel.
            channels_in_group: list[int] = []
            try:
                channels_in_group = sorted({int(c) for c in g_df["spot_channel"].dropna().unique()})
            except Exception:
                channels_in_group = []

            channels_for_title = channels_in_group
            if params.require_spot_channels:
                channels_for_title = sorted({int(c) for c in params.require_spot_channels})

            parts.append(_format_u0_filter(params, channels_for_title))
            if include_desc:
                parts.append(include_desc)

            page_title = "  |  ".join(parts)

            png = render_atlas_page_png(
                batch,
                page_title=page_title,
                params=params,
                image_cache=image_cache,
                manifest_cache=manifest_cache,
                fixed_clim_state=fixed_clim_state,
                context_spots_by_nucleus=context_spots_by_nucleus,
            )

            slide = pres.slides.add_slide(blank)
            bio = BytesIO(png)
            slide.shapes.add_picture(
                bio,
                0,
                0,
                width=pres.slide_width,
                height=pres.slide_height,
            )

    pres.save(str(out_pptx))
    return out_pptx
