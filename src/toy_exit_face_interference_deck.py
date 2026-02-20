"""Generate a small PowerPoint deck illustrating the exit-face interference toy model.

This repo primarily communicates derivations + intuition via Markdown + notebooks.
Occasionally it's useful to have a **slide-ready** artifact that can be pasted into
email threads or vendor discussions.

This module builds a PowerPoint deck (via ``python-pptx``) that:

1) Embeds the conceptual SVG diagram committed under:
   ``docs/figures/speckle_toy_exit_face_interference.svg``

2) Adds a quantitative slide showing how **incoherent averaging over independent
   spectral components** reduces speckle contrast.

The generated ``.pptx`` is intentionally *not* tracked (see ``.gitignore``).
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Sequence

import io


@dataclass(frozen=True)
class ToyExitFaceInterferenceDeckConfig:
    """Configuration for the toy exit-face interference deck."""

    svg_figure_path: Path
    out_pptx_path: Path

    # Slide size (PowerPoint widescreen default: 13.333" x 7.5")
    slide_width_in: float = 13.333
    slide_height_in: float = 7.5

    # Rasterization for the embedded SVG (higher = crisper in PPT)
    svg_raster_width_px: int = 2400

    # Spectral-diversity demo settings
    n_components: Sequence[int] = (1, 5, 20)


def _render_svg_to_png_bytes(*, svg_path: Path, output_width_px: int) -> bytes:
    """Render an SVG file to PNG bytes.

    Notes
    -----
    ``python-pptx`` cannot directly embed SVGs. We keep figures as SVG (repo is
    text-only) and rasterize *at runtime* for slide export.
    """

    if not svg_path.exists():
        raise FileNotFoundError(f"SVG not found: {svg_path}")

    try:
        import cairosvg  # type: ignore
    except Exception as exc:  # pragma: no cover
        raise ImportError(
            "cairosvg is required to rasterize SVG figures for PowerPoint export. "
            "Install it via conda-forge (recommended): `conda install -c conda-forge cairosvg`."
        ) from exc

    png_bytes = cairosvg.svg2png(url=str(svg_path), output_width=int(output_width_px))
    if not isinstance(png_bytes, (bytes, bytearray)):
        raise RuntimeError("cairosvg.svg2png() did not return bytes")
    return bytes(png_bytes)


def _figure_to_png_bytes(*, fig, dpi: int = 180) -> bytes:
    """Serialize a matplotlib Figure to PNG bytes (and close it)."""

    import matplotlib.pyplot as plt

    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=int(dpi), bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf.getvalue()


def _make_spectral_diversity_pdf_plot_png_bytes(*, n_components: Sequence[int]) -> bytes:
    """Create a plot showing the intensity PDF narrowing with N spectral components.

    Model
    -----
    Assume each spectral component produces an independent *fully developed speckle*
    realization at a given pixel.

    Then the *time-averaged / detected* intensity is an incoherent sum:

        I = \\sum_{k=1}^{N} I_k,

    where each I_k is exponential (mean 1 after normalization). The normalized
    intensity u = I/\\langle I \\rangle follows a Gamma distribution with:

        mean(u) = 1,
        std(u)  = 1/\\sqrt{N}.
    """

    # Headless-safe: avoid needing a display when called from scripts/CI.
    import matplotlib

    matplotlib.use("Agg", force=True)

    import matplotlib.pyplot as plt
    import numpy as np
    from scipy.stats import gamma

    n_components = tuple(int(n) for n in n_components)
    if any(n <= 0 for n in n_components):
        raise ValueError("n_components must all be positive")

    # Normalized intensity u = I / mean(I). For Gamma(k=N, theta=1/N): mean=1.
    u = np.linspace(0.0, 3.0, 800)

    fig, ax = plt.subplots(figsize=(10.5, 5.2))

    linestyles = ["-", "--", ":", "-."]
    for i, n in enumerate(n_components):
        k = float(n)
        theta = 1.0 / float(n)
        pdf = gamma(a=k, scale=theta).pdf(u)

        ls = linestyles[i % len(linestyles)]
        ax.plot(u, pdf, color="black", lw=2.0, ls=ls, label=f"N={n}  (C≈{1/np.sqrt(n):.3f})")

    ax.axvline(1.0, color="black", lw=1.2, ls=":")
    ax.text(1.02, ax.get_ylim()[1] * 0.95, "mean", fontsize=11, va="top")

    ax.set_title("Incoherent sum over N independent spectral components narrows the intensity PDF")
    ax.set_xlabel("Normalized intensity  u = I / ⟨I⟩")
    ax.set_ylabel("PDF  p(u)")
    ax.set_xlim(0.0, 3.0)
    ax.grid(True, which="both", lw=0.6, alpha=0.25)

    ax.legend(frameon=False, fontsize=11)

    fig.tight_layout()
    return _figure_to_png_bytes(fig=fig, dpi=200)


def generate_toy_exit_face_interference_deck(*, config: ToyExitFaceInterferenceDeckConfig) -> Path:
    """Generate the PowerPoint deck and write it to disk.

    Parameters
    ----------
    config:
        Deck configuration.

    Returns
    -------
    pathlib.Path
        The output path that was written.
    """

    try:
        from pptx import Presentation  # type: ignore
        from pptx.enum.text import PP_ALIGN  # type: ignore
        from pptx.util import Inches, Pt  # type: ignore
    except Exception as exc:  # pragma: no cover
        raise ImportError(
            "python-pptx is required to generate the toy-model deck. "
            "Install it via environment.yml (conda-forge: python-pptx)."
        ) from exc

    out_path = Path(config.out_pptx_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    pres = Presentation()
    pres.slide_width = Inches(float(config.slide_width_in))
    pres.slide_height = Inches(float(config.slide_height_in))

    blank = pres.slide_layouts[6]

    # --- Slide 1: embed the SVG figure ---
    slide = pres.slides.add_slide(blank)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.25), Inches(12.2), Inches(0.6))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Toy setup: interference at the fiber exit face (one pixel)"
    p.font.size = Pt(28)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    # Figure (rasterized SVG)
    fig_png = _render_svg_to_png_bytes(
        svg_path=Path(config.svg_figure_path),
        output_width_px=int(config.svg_raster_width_px),
    )

    # Keep the figure aspect ratio 900x360 = 2.5.
    # Place it with generous margins.
    fig_left = Inches(0.55)
    fig_top = Inches(1.05)
    fig_width = Inches(12.2)
    fig_height = Inches(12.2 / 2.5)
    slide.shapes.add_picture(io.BytesIO(fig_png), fig_left, fig_top, width=fig_width, height=fig_height)

    # Footer / provenance
    foot = slide.shapes.add_textbox(Inches(0.55), Inches(7.15), Inches(12.2), Inches(0.3))
    ft = foot.text_frame
    ft.clear()
    p2 = ft.paragraphs[0]
    p2.text = f"Source figure: {Path(config.svg_figure_path).as_posix()} (rasterized for PPT)"
    p2.font.size = Pt(11)
    p2.alignment = PP_ALIGN.LEFT

    # --- Slide 2: quantitative spectral-diversity demo ---
    slide2 = pres.slides.add_slide(blank)

    title2 = slide2.shapes.add_textbox(Inches(0.55), Inches(0.25), Inches(12.2), Inches(0.6))
    tf2 = title2.text_frame
    tf2.clear()
    p = tf2.paragraphs[0]
    p.text = "Why wider linewidth reduces speckle: incoherent averaging over spectral components"
    p.font.size = Pt(24)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    plot_png = _make_spectral_diversity_pdf_plot_png_bytes(n_components=config.n_components)
    slide2.shapes.add_picture(
        io.BytesIO(plot_png),
        Inches(0.95),
        Inches(1.25),
        width=Inches(11.5),
    )

    caption = slide2.shapes.add_textbox(Inches(0.95), Inches(6.85), Inches(11.5), Inches(0.6))
    cap_tf = caption.text_frame
    cap_tf.clear()
    pcap = cap_tf.paragraphs[0]
    pcap.text = (
        "If the detector averages out cross-frequency beats, the measured intensity is "
        "I = Σ_k |E_k|² (incoherent). For independent components, speckle contrast scales roughly "
        "C ≈ 1/√N_eff."
    )
    pcap.font.size = Pt(14)
    pcap.alignment = PP_ALIGN.LEFT

    pres.save(str(out_path))
    return out_path
